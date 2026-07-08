"""Stage 1-2: type routing + OCR into normalized page content."""
from __future__ import annotations

import io
import re
import zipfile

from app import ocr as ocr_module
from app.ingestion.rasterize import pdf_page_to_png
from app.ingestion.types import ExtractedContent, PageContent

MIN_CHARS_PER_PAGE = 50  # below this average -> treat PDF as scanned, use OCR

_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff", ".gif"}
_PDF_EXTS = {".pdf"}
_XLSX_EXTS = {".xlsx", ".xlsm"}
_LEGACY_XLS_EXTS = {".xls"}
_CSV_EXTS = {".csv", ".tsv"}
_EMAIL_EXTS = {".eml"}
_TEXT_EXTS = {".txt"}
_ZIP_EXTS = {".zip"}
_DOCX_EXTS = {".docx"}
_LEGACY_DOC_EXTS = {".doc"}
_PPTX_EXTS = {".pptx"}


def _ext(filename: str) -> str:
    idx = filename.rfind(".")
    return filename[idx:].lower() if idx != -1 else ""


# ---------- PDF ----------
def _extract_pdf(data: bytes) -> ExtractedContent:
    import fitz

    doc = fitz.open(stream=data, filetype="pdf")
    try:
        page_texts = [doc[i].get_text("text") for i in range(doc.page_count)]
        page_count = doc.page_count
    finally:
        doc.close()

    total_chars = sum(len(t.strip()) for t in page_texts)
    avg = total_chars / page_count if page_count else 0

    if avg >= MIN_CHARS_PER_PAGE:
        pages = [
            PageContent(page_no=i + 1, text=page_texts[i], ocr_used=False)
            for i in range(page_count)
        ]
        return ExtractedContent(pages=pages, ocr_used=False, kind_detail="pdf_text")

    # Scanned PDF -> rasterize + OCR each page.
    pages = []
    for i in range(page_count):
        png = pdf_page_to_png(data, i)
        page = ocr_module.run_ocr(png)
        pages.append(
            PageContent(
                page_no=i + 1,
                text=page.text,
                ocr_lines=page.lines,
                ocr_used=True,
            )
        )
    return ExtractedContent(pages=pages, ocr_used=True, kind_detail="pdf_ocr")


# ---------- Image ----------
def _extract_image(data: bytes, page_no: int = 1) -> PageContent:
    page = ocr_module.run_ocr(data)
    return PageContent(
        page_no=page_no, text=page.text, ocr_lines=page.lines, ocr_used=True
    )


# ---------- XLSX / CSV ----------
def _extract_xlsx(data: bytes) -> ExtractedContent:
    from openpyxl import load_workbook

    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:
        msg = str(exc).lower()
        if "zip" in msg or "not a zip" in msg:
            raise ValueError(
                "Cannot open spreadsheet: the file may be in legacy .xls format "
                "(not supported by this parser — please re-save as .xlsx) or is "
                "corrupted."
            ) from exc
        raise
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(c.strip() for c in cells):
                parts.append("\t".join(cells))
    wb.close()
    text = "\n".join(parts)
    return ExtractedContent(
        pages=[PageContent(page_no=1, text=text)], kind_detail="xlsx"
    )


def _extract_csv(data: bytes, sep: str = ",") -> ExtractedContent:
    import pandas as pd

    # Auto-detect separator for TSV vs CSV.
    try:
        df = pd.read_csv(
            io.BytesIO(data),
            sep=sep,
            dtype=str,
            keep_default_na=False,
            engine="python",
        )
    except Exception:
        # Last-ditch: decode as plain text so the LLM still sees something.
        text = data.decode("utf-8", errors="replace")
        return ExtractedContent(
            pages=[PageContent(page_no=1, text=text)], kind_detail="csv_raw"
        )
    text = df.to_csv(sep="\t", index=False)
    return ExtractedContent(
        pages=[PageContent(page_no=1, text=text)], kind_detail="csv"
    )


# ---------- Email / text ----------
def _extract_email(data: bytes) -> ExtractedContent:
    import email
    from email import policy

    msg = email.message_from_bytes(data, policy=policy.default)
    body = ""
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                body += part.get_content()
    else:
        body = msg.get_content()
    header = f"From: {msg.get('from', '')}\nSubject: {msg.get('subject', '')}\n\n"
    return ExtractedContent(
        pages=[PageContent(page_no=1, text=header + body)], kind_detail="email"
    )


def _extract_text(data: bytes) -> ExtractedContent:
    text = data.decode("utf-8", errors="replace")
    return ExtractedContent(
        pages=[PageContent(page_no=1, text=text)], kind_detail="text"
    )


# ---------- PowerPoint ----------
def _extract_pptx(data: bytes) -> ExtractedContent:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    pages: list[PageContent] = []
    for i, slide in enumerate(prs.slides):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append("\t".join(cells))
        pages.append(PageContent(page_no=i + 1, text="\n".join(parts)))
    if not pages:
        pages = [PageContent(page_no=1, text="")]
    return ExtractedContent(pages=pages, kind_detail="pptx")


# ---------- Word documents ----------
def _extract_docx(data: bytes) -> ExtractedContent:
    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(data))
    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append("\t".join(cells))
    text = "\n".join(parts)
    return ExtractedContent(
        pages=[PageContent(page_no=1, text=text)], kind_detail="docx"
    )


# ---------- WhatsApp export zip ----------
_WA_LINE_RE = re.compile(
    r"^\[?\d{1,4}[/.\-]\d{1,2}[/.\-]\d{1,4}[,\s]"
)


def parse_whatsapp_chat(chat_text: str) -> str:
    """Keep timestamped messages with sender + date context; drop system noise."""
    lines = []
    for raw in chat_text.replace("\r\n", "\n").split("\n"):
        line = raw.strip()
        if not line:
            continue
        lines.append(line)
    return "\n".join(lines)


def _extract_whatsapp_zip(data: bytes) -> ExtractedContent:
    pages: list[PageContent] = []
    ocr_used = False
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        names = zf.namelist()
        chat_names = [n for n in names if n.lower().endswith("_chat.txt")]
        if not chat_names:
            chat_names = [n for n in names if n.lower().endswith(".txt")]
        chat_text = ""
        for cn in chat_names:
            chat_text += zf.read(cn).decode("utf-8", errors="replace") + "\n"
        pages.append(
            PageContent(page_no=1, text=parse_whatsapp_chat(chat_text), ocr_used=False)
        )

        page_no = 2
        for name in names:
            if _ext(name) in _IMAGE_EXTS:
                try:
                    media = zf.read(name)
                    page = _extract_image(media, page_no=page_no)
                    if page.text.strip():
                        pages.append(page)
                        ocr_used = True
                        page_no += 1
                except Exception:
                    continue
    return ExtractedContent(pages=pages, ocr_used=ocr_used, kind_detail="whatsapp")


# ---------- Dispatcher ----------
def extract_content(filename: str, mime_type: str | None, data: bytes) -> ExtractedContent:
    ext = _ext(filename)
    mime = (mime_type or "").lower()

    if ext in _PDF_EXTS or "pdf" in mime:
        return _extract_pdf(data)
    if ext in _ZIP_EXTS or "zip" in mime:
        return _extract_whatsapp_zip(data)
    if ext in _IMAGE_EXTS or mime.startswith("image/"):
        page = _extract_image(data)
        return ExtractedContent(pages=[page], ocr_used=True, kind_detail="image")
    # Check CSV/TSV by extension FIRST — before MIME-based xlsx routing because
    # some browsers/tools send CSV files with application/vnd.ms-excel MIME type.
    if ext in _CSV_EXTS:
        sep = "\t" if ext == ".tsv" else ","
        return _extract_csv(data, sep=sep)
    if ext in _XLSX_EXTS or "spreadsheet" in mime or "excel" in mime:
        return _extract_xlsx(data)
    if "csv" in mime:
        return _extract_csv(data)
    if ext in _LEGACY_XLS_EXTS:
        raise ValueError(
            "Legacy .xls files are not supported — please re-save as .xlsx and re-upload."
        )
    if ext in _EMAIL_EXTS or "message/rfc822" in mime:
        return _extract_email(data)
    if ext in _DOCX_EXTS or "wordprocessingml" in mime:
        return _extract_docx(data)
    if ext in _PPTX_EXTS or "presentationml" in mime:
        return _extract_pptx(data)
    if ext in _LEGACY_DOC_EXTS or mime == "application/msword":
        raise ValueError(
            "Legacy .doc files are not supported — please re-save as .docx and re-upload."
        )
    if ext in _TEXT_EXTS or mime.startswith("text/"):
        return _extract_text(data)
    # Fallback: try text decode.
    return _extract_text(data)
