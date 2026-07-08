"""M6: global document library ("My Documents") tests. LLM/embeddings mocked."""
from __future__ import annotations

import asyncio

from sqlalchemy import select

from app.db.models import (
    LibraryDocument,
    LibraryDocumentEmbedding,
    ProjectLibraryDocument,
    User,
)
from app.db.session import SessionLocal


def _drain() -> None:
    from app.jobs.worker import drain

    asyncio.run(drain())


def _upload(auth_client, name: str = "past_quote.txt", data: bytes = b"Old solar quote 2024"):
    return auth_client.post(
        "/library/documents",
        files=[("files", (name, data, "text/plain"))],
    )


def test_upload_and_dedup(auth_client):
    r1 = _upload(auth_client)
    assert r1.status_code == 200
    body1 = r1.json()
    assert len(body1["created"]) == 1
    assert body1["skipped"] == []

    # Same bytes again -> deduplicated.
    r2 = _upload(auth_client)
    body2 = r2.json()
    assert body2["created"] == []
    assert len(body2["skipped"]) == 1
    assert body2["skipped"][0]["reason"] == "duplicate"

    docs = auth_client.get("/library/documents").json()
    assert len(docs) == 1


def test_parse_job_populates_embeddings(auth_client):
    r = _upload(auth_client, "notes.txt", b"Transformer installation for HAZECO substation")
    doc_id = r.json()["created"][0]["id"]
    _drain()

    async def _check():
        async with SessionLocal() as session:
            doc = (
                await session.execute(
                    select(LibraryDocument).where(LibraryDocument.id == doc_id)
                )
            ).scalar_one()
            embs = (
                await session.execute(
                    select(LibraryDocumentEmbedding).where(
                        LibraryDocumentEmbedding.library_document_id == doc.id
                    )
                )
            ).scalars().all()
            return doc.status, len(embs)

    status, emb_count = asyncio.run(_check())
    assert status == "parsed"
    assert emb_count >= 1


def test_delete_removes_doc(auth_client):
    r = _upload(auth_client, "del.txt", b"to be deleted")
    doc_id = r.json()["created"][0]["id"]
    _drain()

    resp = auth_client.delete(f"/library/documents/{doc_id}")
    assert resp.status_code == 200
    assert auth_client.get("/library/documents").json() == []


def test_unsupported_type_rejected(auth_client):
    r = auth_client.post(
        "/library/documents",
        files=[("files", ("model.rar", b"\x52\x61\x72\x21", "application/octet-stream"))],
    )
    body = r.json()
    assert body["created"] == []
    assert len(body["errors"]) == 1
    assert "unsupported file type" in body["errors"][0]["error"]
    assert auth_client.get("/library/documents").json() == []


def test_pptx_extraction():
    from io import BytesIO

    from pptx import Presentation
    from pptx.util import Inches

    from app.ingestion.extract import extract_content

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "SMT line quotation summary"
    buf = BytesIO()
    prs.save(buf)

    content = extract_content("deck.pptx", None, buf.getvalue())
    assert "SMT line quotation summary" in content.full_text
    assert content.kind_detail == "pptx"


def test_comments_crud_and_rag(auth_client):
    doc_id = _upload(auth_client, "machine.txt", b"SMT pick and place machine")
    doc_id = doc_id.json()["created"][0]["id"]

    # Add two comments.
    r1 = auth_client.post(
        f"/library/documents/{doc_id}/comments",
        json={"content": "we bought this machine for the H12 lab project"},
    )
    assert r1.status_code == 200
    auth_client.post(
        f"/library/documents/{doc_id}/comments",
        json={"content": "supplier support was slow"},
    )

    comments = auth_client.get(f"/library/documents/{doc_id}/comments").json()
    assert len(comments) == 2

    docs = auth_client.get("/library/documents").json()
    assert docs[0]["comment_count"] == 2

    # RAG: the comment surfaces as a document_comment hit.
    from app.chat.rag import search_all
    from app.db.session import SessionLocal
    from app.llm.embeddings import embed_text
    from sqlalchemy import select as sa_select

    from app.db.models import User

    async def _search():
        async with SessionLocal() as session:
            user = (
                await session.execute(
                    sa_select(User).where(User.email == "user@example.com")
                )
            ).scalar_one()
            hits = await search_all(
                session,
                user.id,
                embed_text("H12 lab project machine"),
                top_k=10,
            )
            return hits

    hits = asyncio.run(_search())
    assert any(h["type"] == "document_comment" for h in hits)

    # Delete one.
    cid = comments[0]["id"]
    r = auth_client.delete(f"/library/documents/{doc_id}/comments/{cid}")
    assert r.status_code == 200
    assert len(auth_client.get(f"/library/documents/{doc_id}/comments").json()) == 1


def test_original_inline_disposition(auth_client):
    doc_id = _upload(auth_client, "view.txt", b"viewable").json()["created"][0]["id"]
    r_dl = auth_client.get(f"/library/documents/{doc_id}/original")
    assert r_dl.headers["content-disposition"].startswith("attachment")
    r_inline = auth_client.get(f"/library/documents/{doc_id}/original?inline=1")
    assert r_inline.headers["content-disposition"].startswith("inline")


def test_link_unlink_project(auth_client):
    pid = auth_client.post("/projects", json={"name": "P"}).json()["id"]
    doc_id = _upload(auth_client, "ref.txt", b"reference doc").json()["created"][0]["id"]

    r = auth_client.post(
        f"/projects/{pid}/library-documents", json={"library_document_id": doc_id}
    )
    assert r.status_code == 200
    assert r.json()["linked"] is True

    # Idempotent second link.
    r2 = auth_client.post(
        f"/projects/{pid}/library-documents", json={"library_document_id": doc_id}
    )
    assert r2.status_code == 200

    linked = auth_client.get(f"/projects/{pid}/library-documents").json()
    assert len(linked) == 1
    assert linked[0]["library_document_id"] == doc_id

    link_id = linked[0]["id"]
    r3 = auth_client.delete(f"/projects/{pid}/library-documents/{link_id}")
    assert r3.status_code == 200
    assert auth_client.get(f"/projects/{pid}/library-documents").json() == []

    # Underlying library doc still exists.
    assert len(auth_client.get("/library/documents").json()) == 1


def test_library_quota_endpoint(auth_client):
    r = auth_client.get("/library/quota")
    assert r.status_code == 200
    body = r.json()
    assert body["limit_bytes"] == 500 * 1024 * 1024
    assert body["used_bytes"] >= 0
    assert body["document_count"] == 0


def test_upload_tracks_size_and_quota(auth_client, monkeypatch):
    from app.config import settings

    monkeypatch.setattr(settings, "library_quota_bytes", 100)

    r = _upload(auth_client, "small.txt", b"x" * 50)
    assert r.status_code == 200
    assert len(r.json()["created"]) == 1

    quota = auth_client.get("/library/quota").json()
    assert quota["used_bytes"] == 50
    assert quota["remaining_bytes"] == 50

    docs = auth_client.get("/library/documents").json()
    assert docs[0]["size_bytes"] == 50

    r2 = _upload(auth_client, "big.txt", b"y" * 60)
    body2 = r2.json()
    assert body2["created"] == []
    assert len(body2["errors"]) == 1
    assert "storage limit exceeded" in body2["errors"][0]["error"]


def test_bulk_delete(auth_client):
    id1 = _upload(auth_client, "a.txt", b"aaa").json()["created"][0]["id"]
    id2 = _upload(auth_client, "b.txt", b"bbb").json()["created"][0]["id"]
    assert len(auth_client.get("/library/documents").json()) == 2

    r = auth_client.post(
        "/library/documents/bulk-delete",
        json={"ids": [id1, id2]},
    )
    assert r.status_code == 200
    assert r.json()["count"] == 2
    assert auth_client.get("/library/documents").json() == []


def test_list_sort_and_project_filter(auth_client):
    pid = auth_client.post("/projects", json={"name": "Alpha"}).json()["id"]
    pid2 = auth_client.post("/projects", json={"name": "Beta"}).json()["id"]

    id_linked = _upload(auth_client, "linked.txt", b"linked doc").json()["created"][0]["id"]
    _upload(auth_client, "solo.txt", b"solo doc")

    auth_client.post(
        f"/projects/{pid}/library-documents",
        json={"library_document_id": id_linked},
    )

    all_docs = auth_client.get("/library/documents").json()
    assert len(all_docs) == 2
    assert all(d["projects"] is not None for d in all_docs)
    linked_doc = next(d for d in all_docs if d["id"] == id_linked)
    assert linked_doc["projects"] == [{"id": pid, "name": "Alpha"}]

    by_project = auth_client.get(f"/library/documents?project_id={pid}").json()
    assert len(by_project) == 1
    assert by_project[0]["id"] == id_linked

    unlinked = auth_client.get("/library/documents?project_id=unlinked").json()
    assert len(unlinked) == 1
    assert unlinked[0]["filename"] == "solo.txt"

    empty = auth_client.get(f"/library/documents?project_id={pid2}").json()
    assert empty == []

    by_name = auth_client.get("/library/documents?sort=name").json()
    assert [d["filename"] for d in by_name] == sorted(
        [d["filename"] for d in by_name]
    )
